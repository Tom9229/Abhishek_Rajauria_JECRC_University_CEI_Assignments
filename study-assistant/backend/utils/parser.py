import fitz # PyMuPDF
import docx
from pptx import Presentation
import os
import re

def clean_text(text: str) -> str:
    # Remove multiple newlines
    text = re.sub(r'\n+', '\n', text)
    # Remove leading/trailing whitespaces per line
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    return "\n".join(lines)

def extract_text_from_file(file_path: str) -> list:
    ext = os.path.splitext(file_path)[1].lower()
    pages = [] # return list of dicts: {"page_content": text, "metadata": {"page": num}}
    
    if ext == ".pdf":
        doc = fitz.open(file_path)
        for i, page in enumerate(doc):
            text = page.get_text()
            text = clean_text(text)
            if text:
                pages.append({"page_content": text, "metadata": {"page": i + 1}})
    elif ext == ".docx":
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        text = clean_text(text)
        if text:
            pages.append({"page_content": text, "metadata": {"page": 1}})
    elif ext == ".pptx":
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text = clean_text(text)
            if text:
                pages.append({"page_content": text, "metadata": {"page": i + 1}})
    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
            text = clean_text(text)
            if text:
                pages.append({"page_content": text, "metadata": {"page": 1}})
    else:
        raise ValueError(f"Unsupported file format: {ext}")
    
    return pages
